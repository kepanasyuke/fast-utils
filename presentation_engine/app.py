"""Web app that turns an uploaded PowerPoint into a branded presentation."""

from __future__ import annotations

import io
import os
import random
import re
import uuid
from pathlib import Path

import qrcode
import requests
import uvicorn
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse, HTMLResponse
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from config import MAX_UPLOAD_BYTES, OUTPUT_DIR as CONFIG_OUTPUT_DIR
from design_system import COLORS, FONTS
from layout_engine import (
    Rect,
    TextBlock,
    choose_font_size,
    split_text_to_fit,
    validate_rect,
)
from source_analysis import analyze_slide, choose_source_template
from utils import create_qr, safe_fit_image


app = FastAPI(title="AI Presentation Brand Engine")
OUTPUT_DIR = CONFIG_OUTPUT_DIR
OUTPUT_DIR.mkdir(exist_ok=True)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
SLIDE_CANVAS = Rect(0, 0, 13.333, 7.5)
WHITE_SPACE_TARGET = 0.70
HEADER_ZONE = (1.0, 0.4, 11.333, 1.0)
HR_ZONE = (1.0, 1.3, 11.333, 0.02)
SAFE_ZONES = {
    "left": (1.0, 1.8, 5.8, 4.5),
    "right": (7.6, 1.8, 4.733, 4.5),
    "center": (1.666, 1.8, 10.0, 4.5),
}
FONT_HEADER = FONTS["header"]
FONT_BODY = FONTS["body"]
FONT_CODE = FONTS["code"]

WHITE = COLORS["white"]
TEXT = COLORS["text"]
MUTED = COLORS["muted"]
LIGHT = COLORS["light"]
BORDER = COLORS["border"]
TERMINAL = COLORS["terminal"]
PALETTES = {
    "cyber_blue": RGBColor(6, 182, 212),
    "matrix_green": RGBColor(16, 185, 129),
    "js_yellow": RGBColor(234, 179, 8),
    "python_blue": RGBColor(59, 130, 246),
    "ruby_red": RGBColor(225, 29, 72),
    "git_orange": RGBColor(249, 115, 22),
    "ai_purple": RGBColor(139, 92, 246),
    "pascal_blue": RGBColor(29, 78, 216),
    "terminal_dark": RGBColor(30, 41, 59),
    "dracula_pink": RGBColor(236, 72, 153),
    "vscode_blue": RGBColor(0, 122, 204),
    "github_gray": RGBColor(36, 41, 46),
}
COLOR_PALETTES = PALETTES

COLOR_WHITE = WHITE
COLOR_TEXT_DARK = TEXT
COLOR_TEXT_MUTED = MUTED
COLOR_CODE_BG = LIGHT
COLOR_BORDER_LIGHT = BORDER
COLOR_TERMINAL_BG = TERMINAL
COLOR_TIMER_RED = RGBColor(239, 68, 68)
COLOR_YANDEX_BLUE = RGBColor(0, 114, 239)

CHARACTER_LIBRARY = [
    {"name": "ИИ-Робот", "url": "https://popsy.co"},
    {"name": "Разработчик", "url": "https://popsy.co"},
    {"name": "Кодер-исследователь", "url": "https://popsy.co"},
    {"name": "Проектировщик ИИ", "url": "https://popsy.co"},
]
RUSSIAN_BUBBLES = [
    "Проверьте входные данные перед запуском алгоритма.",
    "Разбейте сложную задачу на последовательность простых шагов.",
    "Отступы и структура кода важнее случайных исправлений.",
    "Сначала сформулируйте условие, затем выбирайте алгоритм.",
    "Таблица помогает увидеть закономерность быстрее текста.",
    "Сохраняйте промежуточный результат, чтобы проверить ход решения.",
    "Одна понятная схема лучше длинного списка исключений.",
    "Сравните результат с ожидаемым примером.",
    "Названия переменных должны объяснять смысл данных.",
    "Информация становится полезной после правильной обработки.",
]
KEYWORDS_TO_BOLD = [
    "презентац", "слайд", "шаблон", "дизайн", "макет", "гиперссылк", "анимац", "мультимеди",
    "программ", "процессор", "память", "бит", "байт", "код", "алгоритм", "цикл", "переменн",
    "массив", "компиля", "интерпрета", "логик", "истин", "ложь", "таблиц", "граф", "дерево",
    "файл", "папк", "каталог", "вирус", "парол", "защит", "интернет", "сервер", "клиент",
    "браузер", "сайт", "протокол", "систем", "данн", "информац", "процесс", "хранен", "передач",
    "сбор", "обработк", "представлен", "кодирован", "декодирован", "алфавит", "символ", "знак",
    "язык", "текст", "документ", "формат", "абзац", "шрифт", "табуляц", "список", "маркер",
    "изображен", "рисунок", "фотограф", "пиктограмм", "цвет", "контраст", "разрешен", "пиксел",
    "график", "диаграмм", "формул", "величин", "измерен", "количеств", "единиц", "разряд",
    "двоичн", "десятичн", "счислен", "слово", "команд", "оператор", "услови", "ветвлен",
    "функци", "метод", "параметр", "аргумент", "тип", "строк", "числ", "логическ", "ошибк",
    "отладк", "тест", "провер", "результат", "ввод", "вывод", "пользовател", "интерфейс",
    "окно", "кнопк", "меню", "сеть", "адрес", "маршрут", "роутер", "облак", "сервис",
    "сообщен", "электрон", "почт", "ссылк", "поиск", "запрос", "баз", "данных", "запис",
    "поле", "строк", "столбец", "отношен", "безопас", "резерв", "копирован", "носител", "диск",
    "флеш", "оптическ", "магнит", "операцион", "сред", "драйвер", "устройств", "монитор", "клавиатур",
    "мышь", "принтер", "сканер", "микрофон", "видео", "аудио", "мультимедийн", "виртуаль", "цифров",
    "учебник", "задан", "вопрос", "ответ", "практик", "пример", "решен", "этап", "шаг", "последовательн",
    "план", "структур", "сравнен", "определен", "термин", "понят", "свойств", "объект", "явлен", "действ",
    "состояни", "модел", "схем", "соединител", "блок", "условн", "начал", "конец", "стрелк", "узел",
]


class LayoutGridEngine:
    """Контекстный выбор макета без циклического перебора шаблонов."""

    CLASSES = {
        "A": tuple(f"A-{index:02d}" for index in range(1, 13)),
        "B": tuple(f"B-{index:02d}" for index in range(13, 25)),
        "C": tuple(f"C-{index:02d}" for index in range(25, 37)),
        "D": tuple(f"D-{index:02d}" for index in range(37, 49)),
        "E": tuple(f"E-{index:02d}" for index in range(49, 61)),
    }
    SAFE_ZONES = {
        "left": (1.0, 1.8, 5.8, 4.5),
        "right": (7.6, 1.8, 4.733, 4.5),
        "center": (1.666, 1.8, 10.0, 4.5),
    }
    CODE_MARKERS = ("print(", "def ", "for ", "while ", "import ", " = ", "{", "}", "writeln")
    INTERACTIVE_MARKERS = ("?", "задание", "тест", "вопрос", "решите", "найдите", "практика")
    DEFINITION_MARKERS = ("— это", "- это", "называется", "означает", "под термином", "является")
    TIMELINE_MARKERS = ("этап", "шаг", "последовательность", "году", "в 19", "в 20")

    @classmethod
    def analyze_and_score_source(cls, slide_data: dict | str, body_text: str | None = None):
        """Вернуть подробный scoring или публичную метку макета.

        Внешний вызов ``analyze_and_score_source(title, body)`` возвращает
        требуемую метку, внутренний вызов со словарём возвращает diagnostics.
        """

        public_call = isinstance(slide_data, str)
        if public_call:
            slide_data = {"title": slide_data, "body": body_text or ""}

        body = str(slide_data.get("full_body", slide_data.get("body", "")))
        title = str(slide_data.get("title", ""))
        text = f"{title}\n{body}".lower()
        scores = {key: 0 for key in cls.CLASSES}
        evidence = {key: [] for key in cls.CLASSES}

        for marker in cls.CODE_MARKERS:
            if marker in text:
                scores["A"] += 10
                evidence["A"].append(marker)
        for marker in cls.INTERACTIVE_MARKERS:
            if marker in text:
                scores["B"] += 10
                evidence["B"].append(marker)
        for marker in cls.DEFINITION_MARKERS:
            if marker in text:
                scores["C"] += 10
                evidence["C"].append(marker)
        if re.search(r"(?:^|\n)\s*(?:\d+[.)]|[-•])\s+", body) or any(marker in text for marker in cls.TIMELINE_MARKERS):
            scores["D"] += 10
            evidence["D"].append("список/дата/этап")
        if len(body.strip()) < 80 and not scores["A"]:
            scores["E"] += 10
            evidence["E"].append("короткий текст")

        source_profile = slide_data.get("source_profile")
        if source_profile is not None:
            if getattr(source_profile, "table_elements", []):
                scores["D"] += 6
                evidence["D"].append("таблица")
            if getattr(source_profile, "image_elements", []):
                scores["E"] += 4
                evidence["E"].append("изображение")
            if getattr(source_profile, "connector_count", 0) >= 2:
                scores["D"] += 8
                evidence["D"].append("соединители")

        best_score = max(scores.values())
        if best_score == 0:
            winning_class = "E" if len(body) < 180 else "D"
        else:
            winning_class = max(scores, key=lambda key: (scores[key], -ord(key)))

        # Вариация зависит от доказательств и индекса, но не от циклического списка типов.
        fingerprint = sum(ord(char) for char in f"{title}|{body[:240]}")
        variation = (fingerprint + int(slide_data.get("source_index", 0))) % 12
        template_id = cls.CLASSES[winning_class][variation]
        index = int(slide_data.get("source_index", 0))
        result = {
            "scores": scores,
            "evidence": evidence,
            "class": winning_class,
            "template_id": template_id,
            "variation": variation + 1,
            "mode": "mirror-dark" if index % 2 == 0 else "direct-light",
            "safe_zones": cls.SAFE_ZONES,
        }
        if public_call:
            return {
                "A": "code_ide",
                "B": "yandex_quiz",
                "C": "focus_definition",
                "D": "step_timeline",
                "E": "hero_banner",
            }[winning_class]
        return result

    @classmethod
    def apply_semantic_choice(cls, slide_data: dict) -> dict:
        choice = cls.analyze_and_score_source(slide_data)
        slide_data["semantic_layout"] = choice
        return slide_data

    @staticmethod
    def summary_text(text: str, limit: int = 350) -> tuple[str, str]:
        """Вернуть короткий текст слайда и полный текст для заметок."""

        if len(text) <= limit:
            return text, ""
        sentences = [part.strip() for part in re.split(r"(?<=[.!?])\s+", text) if part.strip()]
        summary = " ".join(sentences[:3]).strip()
        if not summary:
            summary = text[:limit].rsplit(" ", 1)[0]
        return summary, text[len(summary):].strip()


def hex_color(color: RGBColor) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def search_and_download_live_gif(query_words: list) -> io.BytesIO | None:
    """Best-effort поиск GIF; сбой сети не должен ломать генерацию PPTX."""

    try:
        query = " ".join(str(word) for word in query_words[:6]) + " animated gif tech"
        response = requests.post(
            "https://duckduckgo.com",
            data={"q": query},
            headers={"User-Agent": "Mozilla/5.0"},
            timeout=4,
        )
        links = re.findall(r"//external-content\.duckduckgo\.com/iu/\?u=[^&]+", response.text)
        candidates = ["https:" + link for link in links if ".gif" in link.lower()]
        if not candidates:
            return None
        image_response = requests.get(random.choice(candidates[:3]), timeout=4)
        image_response.raise_for_status()
        return io.BytesIO(image_response.content)
    except (OSError, requests.RequestException, ValueError):
        return None


def generate_qr_code_stream(link: str) -> io.BytesIO:
    """Создать PNG QR-кода для виджета учебного задания."""

    qr = qrcode.QRCode(version=1, box_size=8, border=2)
    qr.add_data(link or "https://education.yandex.ru/handbook")
    qr.make(fit=True)
    image = qr.make_image(fill_color="black", back_color="white")
    stream = io.BytesIO()
    image.save(stream, format="PNG")
    stream.seek(0)
    return stream


def inject_styled_text_engine(paragraph, text: str, font_name: str, size_pt: int, default_color) -> None:
    """Добавить текст runs с выделением ИТ-терминов по корням слов."""

    paragraph.clear()
    for word in text.split(" "):
        if not word:
            continue
        run = paragraph.add_run()
        run.text = word + " "
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.color.rgb = default_color
        normalized = re.sub(r"[^\wа-яА-Я-ёЁ]", "", word).lower()
        if any(root.lower() in normalized for root in KEYWORDS_TO_BOLD):
            run.font.bold = True


def text_box(slide, x, y, width, height, text, font=FONT_BODY, size=18,
             color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(x, y, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    return shape


def card(slide, x, y, width, height, fill=WHITE, line=BORDER, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    return shape


def qr_stream(link: str) -> io.BytesIO:
    return create_qr(link)


def extract_table(shape) -> list[list[str]]:
    return [
        [re.sub(r"\s+", " ", cell.text).strip() for cell in row.cells]
        for row in shape.table.rows
    ]


def compact_title(value: str) -> str:
    value = re.sub(r"\s+", " ", value).strip()
    return value if len(value) <= 84 else f"{' '.join(value.split()[:10])}..."


def split_content_blocks(lines: list[str]) -> list[str]:
    blocks = []
    for line in lines:
        parts = [part.strip() for part in re.split(r"(?<=[.!?])\s+", line) if part.strip()]
        blocks.extend(parts or [line.strip()])
    return [block for block in blocks if block]


def looks_like_flowchart(text: str, lines: list[str]) -> bool:
    flow_tokens = ("блок-схем", "начало", "конец", "если", "условие", "да", "нет")
    arrows = ("→", "->", "=>", "↓", "шаг")
    short_steps = len(lines) >= 3 and sum(len(line) < 80 for line in lines) >= 3
    return any(token in text for token in flow_tokens) and (short_steps or any(token in text for token in arrows))


def looks_like_timer(text: str) -> bool:
    return bool(re.search(r"\b(?:таймер|время|минут(?:а|ы)?|секунд(?:а|ы)?|\d+\s*(?:мин|сек))\b", text))


def classify_educational_slide(title: str, text: str, blocks: list[str]) -> str | None:
    """Recognize recurring lesson structures before generic heuristics."""

    title_upper = title.upper()
    if "КЛЮЧЕВЫЕ СЛОВА" in title_upper:
        return "key_points"
    if title_upper.startswith(("ВОПРОСЫ", "ЗАДАНИЯ", "ЗАДАЧА")):
        return "exercise_card"
    if title_upper.startswith(("ПРИМЕР", "РЕШЕНИЕ")) or "РЕШЕНИЕ:" in text.upper():
        return "solution_steps"
    if any(token in title_upper for token in ("ФОРМУЛ", "ЗАКОНОМЕРНОСТ", "КОЛИЧЕСТВО")):
        return "formula_focus"
    if any(token in title_upper for token in ("СРАВН", "РАЗЛИЧ", "ВИДЫ")) and len(blocks) >= 2:
        return "comparison_grid"
    if any(token in title_upper for token in ("СХЕМА", "АЛГОРИТМ", "СТРУКТУРА")):
        return "flowchart"
    return None


def prepare_image(blob: bytes, width: int, height: int) -> io.BytesIO:
    return safe_fit_image(blob, width, height)


def iter_shapes(shapes):
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def parse_slides(file_bytes: bytes) -> list[dict[str, str]]:
    try:
        presentation = Presentation(io.BytesIO(file_bytes))
    except Exception as error:
        raise HTTPException(status_code=400, detail="Не удалось прочитать PPTX-файл.") from error

    result = []
    for index, slide in enumerate(presentation.slides):
        source_profile = analyze_slide(slide)
        lines = []
        tables = []
        images = []
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                tables.append(extract_table(shape))
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                images.append(shape.image.blob)
            elif getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    line = re.sub(r"\s+", " ", paragraph.text).strip()
                    if line:
                        prefix = "• " if paragraph.level else ""
                        lines.append(prefix + line)
        if not lines and not tables and not images:
            continue
        title = lines[0] if lines else "Визуальный материал"
        body_lines = lines[1:]

        # Many PowerPoint files keep the title and the entire body in one text box.
        # Split the first sentence so a paragraph cannot become a giant title.
        if len(title) > 90:
            sentence = re.split(r"(?<=[.!?])\s+", title, maxsplit=1)
            if len(sentence) == 2 and len(sentence[0]) >= 18:
                title, first_body = sentence
            else:
                words = title.split()
                title = " ".join(words[:9]) + ("..." if len(words) > 9 else "")
                first_body = " ".join(words[9:])
            if first_body:
                body_lines.insert(0, first_body)
        title = compact_title(title)
        blocks = split_content_blocks(body_lines)
        body = "\n".join(blocks).strip() or title
        display_body, notes_body = LayoutGridEngine.summary_text(body)
        lower = f"{title} {body}".lower()
        educational_type = classify_educational_slide(title, body, blocks)
        source_type = choose_source_template(source_profile, title, body)
        if index == 0 and len(body) < 180 and not tables and not images:
            slide_type = "title_root"
        elif tables:
            slide_type = "table_dashboard"
        elif len(images) >= 2:
            slide_type = "image_collage"
        elif images:
            slide_type = "image_story"
        elif source_type != "auto":
            slide_type = source_type
        elif educational_type:
            slide_type = educational_type
        elif looks_like_flowchart(lower, blocks):
            slide_type = "flowchart"
        elif looks_like_timer(lower):
            slide_type = "timer_focus"
        elif re.search(r"\b(def|class|import|from|return|for|while)\b|[{}]|print\s*\(", lower):
            slide_type = "code_dark_ide" if index % 2 else "code_light_ide"
        elif any(token in lower for token in ("задание", "тест", "вопрос", "практика", "?")):
            slide_type = "yandex_interactive" if index % 2 else "classic_timer"
        elif any(token in lower for token in ("это", "означает", "определение", "термин")) and len(body) < 240:
            slide_type = "big_definition"
        elif any(token in lower for token in ("этап", "история", "шаг", "развитие")):
            slide_type = "step_timeline"
        elif any(token in lower for token in ("git", "репозитор", "коммит", "ветка")):
            slide_type = "git_workflow"
        elif len(body_lines) >= 4:
            slide_type = "key_points"
        elif len(body) < 100:
            slide_type = "minimal_quote" if index % 2 else "accent_statement"
        elif len(body) > 280:
            slide_type = "content_overview"
        else:
            slide_type = ("split_comparison", "process_map", "accent_statement")[index % 3]
        result.append({
            "title": title,
            "body": display_body,
            "full_body": body,
            "notes_body": notes_body,
            "blocks": split_content_blocks(display_body.splitlines()),
            "type": slide_type,
            "table": tables[0] if tables else [],
            "tables": tables,
            "images": images,
            "source_profile": source_profile,
            "source_index": index,
        })
        semantic_choice = LayoutGridEngine.analyze_and_score_source(result[-1])
        result[-1]["semantic_layout"] = semantic_choice
        if not tables and not images and slide_type != "title_root":
            semantic_kind = {
                "A": "code_dark_ide" if index % 2 == 0 else "code_light_ide",
                "B": "exercise_card",
                "C": "big_definition",
                "D": "flowchart",
                "E": "accent_statement",
            }[semantic_choice["class"]]
            result[-1]["type"] = semantic_kind
    return result


def paginate_dataset(dataset: list[dict[str, str]]) -> list[dict[str, str]]:
    paginated = []
    for data in dataset:
        tables = data.get("tables", [])
        if len(tables) > 1:
            for table_index, table in enumerate(tables):
                table_page = dict(data)
                table_page["table"] = table
                table_page["tables"] = [table]
                if table_index:
                    table_page["title"] = f"{data['title']} ({table_index + 1})"
                paginated.append(table_page)
            continue
        blocks = data.get("blocks", [])
        page_size = 6 if data["type"] == "flowchart" else 12 if data["type"] in ("code_dark_ide", "code_light_ide") else 7
        if len(blocks) <= page_size or data["type"] in ("title_root", "image_story", "image_collage", "table_dashboard", "timer_focus"):
            paginated.append(data)
            continue
        measured_pages = split_text_to_fit(
            [TextBlock(text=block) for block in blocks],
            Rect(1.35, 2.85, 10.4, 3.1),
            font_size=14,
            max_blocks=page_size,
        )
        for page_index, page_blocks_objects in enumerate(measured_pages):
            page = dict(data)
            page_blocks = [block.text for block in page_blocks_objects]
            page["blocks"] = page_blocks
            page["body"] = "\n".join(page_blocks)
            if page_index:
                page["title"] = f"{data['title']} ({page_index + 1})"
            if data["type"] == "content_overview":
                page["type"] = "key_points"
            paginated.append(page)
    return paginated


def add_header(slide, title: str, accent: RGBColor) -> None:
    header_zone = Rect(1.0, 0.4, 11.333, 1.0)
    title_size = 22 if len(title) > 45 else 32
    title_size = choose_font_size(title, header_zone, title_size, minimum=18)
    text_box(slide, Inches(1.0), Inches(0.4), Inches(11.333), Inches(1.0),
             title, font=FONT_HEADER, size=title_size, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.3), Inches(1.0), Inches(0.02))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def add_speaker_notes(slide, data: dict[str, str]) -> None:
    """Store omitted source text in speaker notes without cluttering the slide."""

    notes_body = data.get("notes_body", "").strip()
    if not notes_body:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "Подробный материал исходного слайда:\n\n" + notes_body


def add_table_visual(slide, values: list[list[str]], accent: RGBColor) -> None:
    if not values:
        return
    rows = len(values)
    columns = max(len(row) for row in values)
    table_shape = slide.shapes.add_table(
        rows, columns, Inches(0.95), Inches(1.8), Inches(11.45), Inches(4.85)
    )
    table = table_shape.table
    for column in range(columns):
        table.columns[column].width = Inches(11.45 / columns)
    for row_index, row in enumerate(values):
        for column_index in range(columns):
            cell = table.cell(row_index, column_index)
            cell.text = row[column_index] if column_index < len(row) else ""
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)
            cell.fill.solid()
            cell.fill.fore_color.rgb = accent if row_index == 0 else (LIGHT if row_index % 2 else WHITE)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT_HEADER if row_index == 0 else FONT_BODY
                paragraph.font.size = Pt(13 if row_index == 0 else 12)
                paragraph.font.bold = row_index == 0
                paragraph.font.color.rgb = WHITE if row_index == 0 else TEXT
                paragraph.alignment = PP_ALIGN.LEFT


def add_picture_frame(slide, image_blob: bytes, x, y, width, height, accent: RGBColor) -> None:
    frame = card(slide, x - Inches(0.08), y - Inches(0.08), width + Inches(0.16), height + Inches(0.16), WHITE, accent)
    frame.shadow.inherit = False
    pixel_width = max(1, int(width / Inches(1) * 150))
    pixel_height = max(1, int(height / Inches(1) * 150))
    slide.shapes.add_picture(prepare_image(image_blob, pixel_width, pixel_height), x, y, width=width, height=height)


def add_image_placeholder(slide, x, y, width, height, accent: RGBColor) -> None:
    placeholder = card(slide, x, y, width, height, LIGHT, accent)
    text_frame = placeholder.text_frame
    text_frame.text = "ВИЗУАЛЬНЫЙ АКЦЕНТ"
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = FONT_HEADER
    paragraph.font.size = Pt(13)
    paragraph.font.bold = True
    paragraph.font.color.rgb = accent


def add_editorial_image_layout(slide, data: dict[str, str], accent: RGBColor) -> None:
    images = data.get("images", [])
    primary_image = images[0] if images else None
    image_x, image_y = Inches(0.95), Inches(1.72)
    image_width, image_height = Inches(6.15), Inches(4.85)
    if primary_image:
        add_picture_frame(slide, primary_image, image_x, image_y, image_width, image_height, accent)
    else:
        add_image_placeholder(slide, image_x, image_y, image_width, image_height, accent)

    text_card = card(slide, Inches(7.55), Inches(1.72), Inches(4.8), Inches(4.85), LIGHT, BORDER)
    text_card.line.color.rgb = BORDER
    body = data.get("body", "").strip()
    body_zone = Rect(7.95, 2.12, 4.0, 3.95)
    body_size = choose_font_size(body, body_zone, 18, minimum=11)
    text_box(slide, Inches(7.95), Inches(2.12), Inches(4.0), Inches(3.95), body,
             size=body_size, color=MUTED)

    # Extra images become a small visual index below the main composition.
    for image_index, image in enumerate(images[1:4]):
        thumb_x = Inches(0.95 + image_index * 1.5)
        add_picture_frame(slide, image, thumb_x, Inches(6.72), Inches(1.25), Inches(0.58), accent)


def add_flowchart_layout(slide, data: dict[str, str], accent: RGBColor) -> None:
    blocks = data.get("blocks", [])[:6]
    if not blocks:
        blocks = [data.get("body", "")]
    if len(blocks) > 4:
        columns = 2
        rows = (len(blocks) + columns - 1) // columns
        node_width = Inches(5.25)
        node_height = Inches(1.05)
        x_positions = (Inches(0.95), Inches(6.95))
        y_step = 4.55 / max(1, rows)
        for node_index, block in enumerate(blocks):
            column = node_index % columns
            row = node_index // columns
            x = x_positions[column]
            y = Inches(1.65 + row * y_step)
            shape_type = MSO_SHAPE.DIAMOND if any(token in block.lower() for token in ("если", "условие", "?")) else MSO_SHAPE.ROUNDED_RECTANGLE
            node = slide.shapes.add_shape(shape_type, x, y, node_width, node_height)
            node.fill.solid(); node.fill.fore_color.rgb = LIGHT if node_index % 2 else WHITE
            node.line.color.rgb = accent; node.line.width = Pt(1.8)
            text_box(slide, x + Inches(0.25), y + Inches(0.2), node_width - Inches(0.5), Inches(0.58),
                     block, size=12 if len(block) > 55 else 14, color=TEXT,
                     bold=node_index == 0, align=PP_ALIGN.CENTER)
            if row < rows - 1:
                connector = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x + Inches(2.35), y + node_height + Inches(0.05), Inches(0.5), Inches(0.22))
                connector.fill.solid(); connector.fill.fore_color.rgb = accent; connector.line.fill.background()
        return
    start_x = Inches(1.0)
    node_width = Inches(3.3)
    node_height = Inches(0.72)
    gap = Inches(0.34)
    for node_index, block in enumerate(blocks):
        y = Inches(1.75) + node_index * (node_height + gap)
        shape_type = MSO_SHAPE.DIAMOND if any(token in block.lower() for token in ("если", "условие", "?") ) else MSO_SHAPE.ROUNDED_RECTANGLE
        node = slide.shapes.add_shape(shape_type, start_x, y, node_width, node_height)
        node.fill.solid()
        node.fill.fore_color.rgb = LIGHT if node_index % 2 else WHITE
        node.line.color.rgb = accent
        node.line.width = Pt(1.8)
        text_box(slide, start_x + Inches(0.2), y + Inches(0.12), node_width - Inches(0.4), Inches(0.42),
                 block, size=11 if len(block) > 42 else 13, color=TEXT, bold=node_index == 0,
                 align=PP_ALIGN.CENTER)
        if node_index < len(blocks) - 1:
            connector = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.42), y + node_height + Inches(0.03), Inches(0.45), Inches(0.24))
            connector.fill.solid(); connector.fill.fore_color.rgb = accent; connector.line.fill.background()
    card(slide, Inches(5.45), Inches(1.8), Inches(6.85), Inches(4.6), LIGHT, BORDER)
    text_box(slide, Inches(5.9), Inches(2.25), Inches(5.95), Inches(0.45), "ЛОГИКА АЛГОРИТМА",
             font=FONT_HEADER, size=15, color=accent, bold=True)
    text_box(slide, Inches(5.9), Inches(3.0), Inches(5.7), Inches(2.2),
             "Последовательность действий собрана в понятную схему: каждый шаг имеет собственную зону и не смешивается с соседним текстом.",
             size=16, color=MUTED)


def add_timer_layout(slide, data: dict[str, str], accent: RGBColor) -> None:
    match = re.search(r"\b\d+\s*(?:минут(?:а|ы)?|мин|секунд(?:а|ы)?|сек)\b", data.get("body", ""), re.IGNORECASE)
    timer_value = match.group(0) if match else "02 мин"
    card(slide, Inches(0.95), Inches(1.75), Inches(4.3), Inches(4.75), TERMINAL, accent)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.1), Inches(2.25), Inches(1.95), Inches(1.95))
    circle.fill.background(); circle.line.color.rgb = accent; circle.line.width = Pt(3)
    text_box(slide, Inches(1.45), Inches(4.45), Inches(3.3), Inches(0.6), timer_value.upper(),
             font=FONT_HEADER, size=27, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, Inches(1.4), Inches(5.35), Inches(3.4), Inches(0.45), "ЛИМИТ НА РЕШЕНИЕ",
             font=FONT_HEADER, size=12, color=accent, bold=True, align=PP_ALIGN.CENTER)
    card(slide, Inches(5.75), Inches(1.75), Inches(6.6), Inches(4.75), LIGHT, BORDER)
    text_box(slide, Inches(6.2), Inches(2.2), Inches(5.7), Inches(0.45), "ЗАДАНИЕ",
             font=FONT_HEADER, size=15, color=accent, bold=True)
    text_box(slide, Inches(6.2), Inches(2.9), Inches(5.7), Inches(2.7), data.get("body", ""),
             size=17 if len(data.get("body", "")) < 260 else 13, color=MUTED)


def add_code_layout(slide, data: dict[str, str], accent: RGBColor, dark: bool) -> None:
    fill = TERMINAL if dark else LIGHT
    color = WHITE if dark else TEXT
    blocks = data.get("blocks", []) or data.get("body", "").splitlines()
    code_lines = blocks[:16]
    code = "\n".join(code_lines)
    card(slide, Inches(0.9), Inches(1.7), Inches(7.0), Inches(4.95), fill, accent, False)
    text_box(slide, Inches(1.2), Inches(2.02), Inches(0.45), Inches(4.25),
             "\n".join(str(index + 1) for index in range(len(code_lines))),
             font=FONT_CODE, size=11, color=accent, align=PP_ALIGN.RIGHT)
    text_box(slide, Inches(1.85), Inches(2.02), Inches(5.65), Inches(4.25), code,
             font=FONT_CODE, size=12 if len(code_lines) > 9 else 14, color=color)
    card(slide, Inches(8.45), Inches(1.7), Inches(3.9), Inches(4.95), WHITE, accent)
    text_box(slide, Inches(8.85), Inches(2.2), Inches(3.1), Inches(0.5), "КОДОВЫЙ БЛОК",
             font=FONT_HEADER, size=15, color=accent, bold=True)
    text_box(slide, Inches(8.85), Inches(3.0), Inches(3.0), Inches(2.3),
             "Синтаксис отделён от пояснения. Номера строк помогают обсуждать конкретный шаг алгоритма.",
             size=15, color=MUTED)


def add_exercise_layout(slide, data: dict[str, str], accent: RGBColor, index: int) -> None:
    """Render a calm worksheet instead of forcing every task into a timer."""

    card(slide, Inches(0.95), Inches(1.75), Inches(8.15), Inches(4.85), LIGHT, BORDER)
    text_box(slide, Inches(1.35), Inches(2.15), Inches(7.35), Inches(0.42), "РАБОЧИЙ ЛИСТ",
             font=FONT_HEADER, size=15, color=accent, bold=True)
    exercise_zone = Rect(1.35, 2.85, 7.25, 3.2)
    font_size = choose_font_size(data.get("body", ""), exercise_zone, 19, minimum=12)
    text_box(slide, Inches(1.35), Inches(2.85), Inches(7.25), Inches(3.2), data.get("body", ""),
             size=font_size, color=TEXT)
    card(slide, Inches(9.55), Inches(1.75), Inches(2.75), Inches(4.85), WHITE, accent)
    text_box(slide, Inches(9.9), Inches(2.2), Inches(2.05), Inches(0.6), f"ЗАДАНИЕ {index + 1:02d}",
             font=FONT_HEADER, size=14, color=accent, bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, Inches(9.9), Inches(3.15), Inches(2.05), Inches(1.5),
             "Сформулируйте ответ\nи обсудите ход решения.", size=15, color=MUTED, align=PP_ALIGN.CENTER)
    timer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.82), Inches(5.1), Inches(2.2), Inches(0.72))
    timer.fill.solid(); timer.fill.fore_color.rgb = TERMINAL; timer.line.color.rgb = accent
    text_box(slide, Inches(9.98), Inches(5.25), Inches(1.9), Inches(0.35), "ВРЕМЯ ПОШЛО · 02:00",
             font=FONT_HEADER, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    slide.shapes.add_picture(qr_stream(data.get("link", "")), Inches(10.42), Inches(1.95), Inches(1.0), Inches(1.0))


def add_solution_layout(slide, data: dict[str, str], accent: RGBColor) -> None:
    steps = data.get("blocks", [])[:6] or [data.get("body", "")]
    for step_index, step in enumerate(steps):
        x = Inches(0.95 + (step_index % 2) * 6.0)
        y = Inches(1.75 + (step_index // 2) * 1.55)
        card(slide, x, y, Inches(5.25), Inches(1.12), LIGHT if step_index % 2 else WHITE, accent)
        text_box(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.55), Inches(0.5),
                 f"{step_index + 1:02d}", font=FONT_HEADER, size=17, color=accent, bold=True)
        text_box(slide, x + Inches(1.0), y + Inches(0.17), Inches(3.95), Inches(0.72), step,
                 size=12 if len(step) > 110 else 14, color=TEXT)


def add_formula_layout(slide, data: dict[str, str], accent: RGBColor) -> None:
    card(slide, Inches(0.95), Inches(1.75), Inches(11.45), Inches(2.25), TERMINAL, accent)
    formula = data.get("blocks", [data.get("body", "")])[0]
    text_box(slide, Inches(1.35), Inches(2.4), Inches(10.65), Inches(0.9), formula,
             font=FONT_CODE, size=26 if len(formula) < 45 else 18, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    text_box(slide, Inches(1.35), Inches(4.55), Inches(10.6), Inches(1.2),
             "Разберите обозначения и объясните, как формула связана с задачей.",
             size=18, color=MUTED, align=PP_ALIGN.CENTER)


def add_comparison_layout(slide, data: dict[str, str], accent: RGBColor) -> None:
    blocks = data.get("blocks", [])
    split = max(1, len(blocks) // 2)
    columns = [blocks[:split], blocks[split:]]
    labels = ("ПЕРВЫЙ ВАРИАНТ", "ВТОРОЙ ВАРИАНТ")
    for column_index, column in enumerate(columns):
        x = Inches(0.95 + column_index * 6.0)
        card(slide, x, Inches(1.75), Inches(5.25), Inches(4.9), LIGHT if column_index else WHITE, accent)
        text_box(slide, x + Inches(0.35), Inches(2.15), Inches(4.55), Inches(0.45), labels[column_index],
                 font=FONT_HEADER, size=14, color=accent, bold=True)
        text_box(slide, x + Inches(0.35), Inches(2.9), Inches(4.55), Inches(3.1),
                 "\n".join(f"• {item}" for item in column) or data.get("body", ""),
                 size=14, color=TEXT)


def render_slide(slide, data: dict[str, str], accent: RGBColor, link: str, index: int) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    kind = data["type"]
    title, body = data["title"], data["body"]

    if kind == "title_root":
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.45), SLIDE_HEIGHT)
        stripe.fill.solid(); stripe.fill.fore_color.rgb = accent; stripe.line.fill.background()
        card(slide, Inches(1.4), Inches(1.55), Inches(10.5), Inches(4.45), LIGHT, BORDER)
        title_size = 42 if len(title) < 32 else 34 if len(title) < 58 else 27
        text_box(slide, Inches(1.8), Inches(2.15), Inches(9.7), Inches(1.55), title,
             font=FONT_HEADER, size=title_size, bold=True, align=PP_ALIGN.CENTER)
        text_box(slide, Inches(2.1), Inches(3.95), Inches(9.1), Inches(1.45), body,
             size=16 if len(body) < 120 else 13, color=MUTED, align=PP_ALIGN.CENTER)
        return

    add_header(slide, title, accent)
    if kind == "table_dashboard":
        text_box(slide, Inches(0.95), Inches(1.42), Inches(11.45), Inches(0.28),
                 "Данные собраны в единую визуальную систему", size=11, color=MUTED)
        add_table_visual(slide, data.get("table", []), accent)
    elif kind in ("image_collage", "image_story"):
        add_editorial_image_layout(slide, data, accent)
    elif kind == "flowchart":
        add_flowchart_layout(slide, data, accent)
    elif kind == "timer_focus":
        add_timer_layout(slide, data, accent)
    elif kind == "exercise_card":
        add_exercise_layout(slide, data, accent, index)
    elif kind == "solution_steps":
        add_solution_layout(slide, data, accent)
    elif kind == "formula_focus":
        add_formula_layout(slide, data, accent)
    elif kind == "comparison_grid":
        add_comparison_layout(slide, data, accent)
    elif kind == "key_points":
        card(slide, Inches(0.95), Inches(1.75), Inches(11.45), Inches(4.9), LIGHT, BORDER)
        points = [line.strip("• ") for line in data.get("blocks", []) if line.strip()]
        for point_index, point in enumerate(points[:7]):
            y = 2.15 + point_index * 0.6
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), Inches(y + 0.04), Inches(0.24), Inches(0.24))
            marker.fill.solid(); marker.fill.fore_color.rgb = accent; marker.line.fill.background()
            text_box(slide, Inches(1.8), Inches(y), Inches(10.0), Inches(0.48), point,
                     size=15 if len(point) < 100 else 12, color=TEXT)
    elif kind == "accent_statement":
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(1.8), Inches(0.28), Inches(4.8))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = accent; stripe.line.fill.background()
        text_box(slide, Inches(1.7), Inches(2.15), Inches(10.2), Inches(0.45), "ГЛАВНАЯ МЫСЛЬ",
                 font=FONT_HEADER, size=15, color=accent, bold=True)
        text_box(slide, Inches(1.7), Inches(2.85), Inches(10.1), Inches(2.4), body,
                 font=FONT_HEADER, size=30 if len(body) < 140 else 21,
                 color=TEXT, bold=True)
    elif kind == "process_map":
        steps = [part.strip(" •") for part in re.split(r"\n|(?<=[.!?])\s+", body) if part.strip()][:4]
        for step_index, step in enumerate(steps):
            x = Inches(0.95 + (step_index % 2) * 6.0)
            y = Inches(1.85 + (step_index // 2) * 2.45)
            card(slide, x, y, Inches(5.25), Inches(1.75), LIGHT if step_index % 2 else WHITE, accent)
            text_box(slide, x + Inches(0.28), y + Inches(0.22), Inches(0.55), Inches(0.4),
                     f"{step_index + 1:02d}", font=FONT_HEADER, size=18, color=accent, bold=True)
            text_box(slide, x + Inches(1.0), y + Inches(0.24), Inches(3.9), Inches(1.2), step,
                     size=14 if len(step) < 110 else 11, color=TEXT)
    elif kind == "content_overview":
        card(slide, Inches(0.95), Inches(1.75), Inches(11.45), Inches(4.9), LIGHT, BORDER)
        text_box(slide, Inches(1.35), Inches(2.2), Inches(10.5), Inches(0.4), "ОБЗОР МАТЕРИАЛА",
                 font=FONT_HEADER, size=15, color=accent, bold=True)
        overview = "\n\n".join(f"• {block}" for block in data.get("blocks", [])[:7]) or body
        text_box(slide, Inches(1.35), Inches(2.85), Inches(10.4), Inches(3.1), overview,
                 size=13 if len(overview) < 650 else 10, color=TEXT)
    elif kind in ("code_dark_ide", "code_light_ide"):
        add_code_layout(slide, data, accent, kind == "code_dark_ide")
    elif kind in ("yandex_interactive", "classic_timer"):
        card(slide, Inches(0.9), Inches(1.8), Inches(7.4), Inches(4.7), WHITE, accent)
        text_box(slide, Inches(1.35), Inches(2.25), Inches(6.5), Inches(2.5), body, size=19, color=MUTED)
        label = "ЗАДАНИЕ В ЯНДЕКС УЧЕБНИКЕ" if kind == "yandex_interactive" else "ВРЕМЯ ПОШЛО"
        text_box(slide, Inches(1.35), Inches(5.2), Inches(5.2), Inches(0.5), label,
                 font=FONT_HEADER, size=16, color=accent, bold=True)
        slide.shapes.add_picture(qr_stream(link), Inches(10.0), Inches(2.2), Inches(1.8), Inches(1.8))
    elif kind == "big_definition":
        card(slide, Inches(1.0), Inches(2.0), Inches(11.2), Inches(2.8), WHITE, accent)
        text_box(slide, Inches(1.5), Inches(2.5), Inches(10.2), Inches(0.5), "КЛЮЧЕВОЙ ТЕРМИН",
                 font=FONT_HEADER, size=15, color=accent, bold=True)
        text_box(slide, Inches(1.5), Inches(3.15), Inches(10.2), Inches(1.1), body,
                 font=FONT_HEADER, size=24 if len(body) < 100 else 18, bold=True)
    elif kind == "step_timeline":
        steps = [part.strip() for part in re.split(r"[.!?]\s*", body) if part.strip()][:3] or [body]
        for step_index, step in enumerate(steps):
            x = Inches(1.1 + step_index * 4.0)
            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(3.0), Inches(0.55), Inches(0.55))
            node.fill.solid(); node.fill.fore_color.rgb = accent; node.line.fill.background()
            text_box(slide, x - Inches(0.2), Inches(2.25), Inches(2.7), Inches(0.5),
                     f"ШАГ {step_index + 1}", font=FONT_HEADER, size=14, color=accent, bold=True)
            text_box(slide, x - Inches(0.2), Inches(3.8), Inches(3.2), Inches(1.4), step, size=15, color=MUTED)
            if step_index < len(steps) - 1:
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.55), Inches(3.24), Inches(3.4), Inches(0.06))
                line.fill.solid(); line.fill.fore_color.rgb = accent; line.line.fill.background()
    elif kind == "git_workflow":
        card(slide, Inches(1.0), Inches(1.9), Inches(11.2), Inches(4.5), LIGHT, accent, False)
        text_box(slide, Inches(1.5), Inches(2.4), Inches(10.0), Inches(0.5), "GIT WORKFLOW",
                 font=FONT_HEADER, size=18, color=accent, bold=True)
        text_box(slide, Inches(1.5), Inches(3.15), Inches(10.0), Inches(2.2), body,
                 font=FONT_CODE, size=17)
    elif kind == "minimal_quote":
        text_box(slide, Inches(1.0), Inches(2.2), Inches(11.0), Inches(2.8), f'« {body} »',
                 font=FONT_HEADER, size=30, color=accent, bold=True, align=PP_ALIGN.CENTER)
    else:
        card(slide, Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.9), LIGHT, accent)
        text_box(slide, Inches(1.35), Inches(2.25), Inches(10.5), Inches(0.5), "КОНТЕКСТ", font=FONT_HEADER, size=15, color=accent, bold=True)
        body_size = 17 if len(body) < 300 else 14
        text_box(slide, Inches(1.35), Inches(2.95), Inches(10.5), Inches(3.2), body,
             size=body_size, color=MUTED)


def validate_slide_geometry(slide) -> list[str]:
    """Report generated shapes that exceed the fixed 16:9 slide canvas."""

    issues = []
    for shape_index, shape in enumerate(slide.shapes):
        rect = Rect(
            shape.left / Inches(1),
            shape.top / Inches(1),
            shape.width / Inches(1),
            shape.height / Inches(1),
        )
        for issue in validate_rect(rect, SLIDE_CANVAS):
            issues.append(f"shape {shape_index}: {issue.message}")
    return issues


def build_presentation(dataset: list[dict[str, str]], palette: str, link: str) -> Path:
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    accent = PALETTES.get(palette, PALETTES["cyber_blue"])
    blank = presentation.slide_layouts[6]
    layout_issues = []
    for index, data in enumerate(paginate_dataset(dataset)):
        data = dict(data)
        data["link"] = link
        slide = presentation.slides.add_slide(blank)
        render_slide(slide, data, accent, link, index)
        add_speaker_notes(slide, data)
        layout_issues.extend((index + 1, issue) for issue in validate_slide_geometry(slide))
    if layout_issues:
        issue_preview = "; ".join(f"слайд {index}: {message}" for index, message in layout_issues[:5])
        print(f"Предупреждения компоновки: {issue_preview}")
    output = OUTPUT_DIR / f"generated_{uuid.uuid4().hex[:10]}.pptx"
    presentation.save(output)
    return output


def build_presentation_industrial(dataset: list, chosen_palette: str, yandex_link: str) -> str:
    """Совместимый публичный API промышленной сборки презентации."""

    return str(build_presentation(dataset, chosen_palette, yandex_link))


PALETTE_OPTIONS = "".join(
    f'<option value="{key}">{key.replace("_", " ").title()}</option>'
    for key in PALETTES
)
HTML_UI = f"""<!doctype html>
<html lang="ru"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>Presentation Brand Engine</title><style>
body{{margin:0;padding:40px;background:#020617;color:#f8fafc;font-family:'Segoe UI',sans-serif}}
main{{max-width:760px;margin:auto;padding:42px;background:#0f172a;border:1px solid #1e293b;border-radius:24px;box-shadow:0 24px 70px #0008}}
h1{{margin:0 0 12px;color:#38bdf8;font-size:34px}}p{{color:#94a3b8;line-height:1.6}}
label{{display:block;margin:24px 0 8px;color:#cbd5e1;font-weight:700;font-size:13px;text-transform:uppercase}}
input,select,button{{width:100%;box-sizing:border-box;padding:14px;border-radius:10px;font:inherit}}
input,select{{background:#020617;color:#fff;border:1px solid #334155}}button{{margin-top:28px;border:0;background:#38bdf8;color:#082f49;font-weight:800;cursor:pointer}}
button:hover{{background:#7dd3fc}}small{{color:#64748b}}
</style></head><body><main><h1>AI Presentation Brand Engine</h1>
<p>Загрузите презентацию, выберите акцент и скачайте новую версию с единым визуальным стилем.</p>
<form action="/engine-process" method="post" enctype="multipart/form-data">
<label for="file">Исходная презентация</label><input id="file" name="file" type="file" accept=".pptx" required>
<label for="palette">Палитра</label><select id="palette" name="palette">{PALETTE_OPTIONS}</select>
<label for="yandex_link">Ссылка для QR-кода (необязательно)</label><input id="yandex_link" name="yandex_link" type="url" placeholder="https://...">
<button type="submit">Собрать редизайн</button></form><small>Поддерживается формат .pptx</small></main></body></html>"""


@app.get("/", response_class=HTMLResponse)
async def root() -> str:
    return HTML_UI


@app.post("/engine-process")
async def engine_process(
    palette: str = Form("cyber_blue"),
    yandex_link: str = Form(""),
    file: UploadFile = File(...),
):
    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Загрузите файл в формате .pptx.")
    data = await file.read()
    if len(data) > MAX_UPLOAD_BYTES:
        raise HTTPException(status_code=413, detail="Файл презентации слишком большой.")
    slides = parse_slides(data)
    if not slides:
        raise HTTPException(status_code=400, detail="В презентации не найден текстовый контент.")
    output = build_presentation(slides, palette, yandex_link)
    return FileResponse(output, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        filename=f"redesigned_{Path(file.filename).stem}.pptx")


if __name__ == "__main__":
    uvicorn.run("app:app", host=os.getenv("HOST", "0.0.0.0"),
                port=int(os.getenv("PORT", "8000")), reload=True)