"""Source-slide analysis used to choose stable redesign templates.

The analyzer intentionally keeps only layout facts that are useful to a redesign:
object kind, normalized bounds, text density, image aspect ratio, and likely role.
It does not try to clone the original slide pixel-for-pixel.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from enum import Enum
from statistics import mean
from typing import Iterable

from pptx.enum.shapes import MSO_SHAPE_TYPE


class ElementKind(str, Enum):
    TEXT = "text"
    IMAGE = "image"
    TABLE = "table"
    CONNECTOR = "connector"
    SHAPE = "shape"


class TextRole(str, Enum):
    TITLE = "title"
    SUBTITLE = "subtitle"
    BODY = "body"
    LABEL = "label"
    UNKNOWN = "unknown"


@dataclass(frozen=True)
class Bounds:
    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def area(self) -> float:
        return max(0.0, self.width) * max(0.0, self.height)

    @property
    def center_x(self) -> float:
        return self.left + self.width / 2

    @property
    def center_y(self) -> float:
        return self.top + self.height / 2

    @property
    def aspect_ratio(self) -> float:
        return self.width / self.height if self.height else 0.0

    def normalized(self, slide_width: float, slide_height: float) -> "Bounds":
        return Bounds(
            self.left / slide_width,
            self.top / slide_height,
            self.width / slide_width,
            self.height / slide_height,
        )


@dataclass
class SourceElement:
    kind: ElementKind
    bounds: Bounds
    text: str = ""
    role: TextRole = TextRole.UNKNOWN
    paragraph_count: int = 0
    row_count: int = 0
    column_count: int = 0
    has_bullets: bool = False
    image_blob: bytes | None = None
    z_order: int = 0

    @property
    def text_length(self) -> int:
        return len(self.text)


@dataclass
class SourceProfile:
    slide_width: float
    slide_height: float
    elements: list[SourceElement] = field(default_factory=list)

    @property
    def text_elements(self) -> list[SourceElement]:
        return [element for element in self.elements if element.kind == ElementKind.TEXT]

    @property
    def image_elements(self) -> list[SourceElement]:
        return [element for element in self.elements if element.kind == ElementKind.IMAGE]

    @property
    def table_elements(self) -> list[SourceElement]:
        return [element for element in self.elements if element.kind == ElementKind.TABLE]

    @property
    def connector_count(self) -> int:
        return sum(element.kind == ElementKind.CONNECTOR for element in self.elements)

    @property
    def text_density(self) -> float:
        text_area = sum(element.bounds.area for element in self.text_elements)
        canvas_area = self.slide_width * self.slide_height
        return text_area / canvas_area if canvas_area else 0.0

    @property
    def visual_density(self) -> float:
        visual_area = sum(
            element.bounds.area
            for element in self.elements
            if element.kind in (ElementKind.IMAGE, ElementKind.TABLE)
        )
        canvas_area = self.slide_width * self.slide_height
        return visual_area / canvas_area if canvas_area else 0.0

    @property
    def has_two_columns(self) -> bool:
        text = self.text_elements
        if len(text) < 2:
            return False
        left = sum(element.bounds.center_x < self.slide_width / 2 for element in text)
        right = len(text) - left
        return left > 0 and right > 0

    @property
    def dominant_image_ratio(self) -> float:
        ratios = [element.bounds.aspect_ratio for element in self.image_elements if element.bounds.aspect_ratio]
        return mean(ratios) if ratios else 0.0

    def dominant_text_role(self) -> TextRole:
        roles = [element.role for element in self.text_elements]
        if TextRole.TITLE in roles:
            return TextRole.TITLE
        if TextRole.SUBTITLE in roles:
            return TextRole.SUBTITLE
        return TextRole.BODY


def _bounds(shape, slide_width: float, slide_height: float) -> Bounds:
    return Bounds(
        shape.left / 914400,
        shape.top / 914400,
        shape.width / 914400,
        shape.height / 914400,
    )


def _text_role(shape, text: str, slide_width: float, slide_height: float) -> TextRole:
    bounds = _bounds(shape, slide_width, slide_height)
    upper = text.strip().upper()
    if bounds.top < slide_height * 0.22 and len(text) < 120:
        return TextRole.TITLE
    if len(text) < 90 and bounds.height < 0.75:
        return TextRole.LABEL
    if upper.startswith(("ПРИМЕР", "ЗАДАНИЕ", "ВОПРОС", "РЕШЕНИЕ")):
        return TextRole.SUBTITLE
    return TextRole.BODY


def iter_source_shapes(shapes) -> Iterable:
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_source_shapes(shape.shapes)
        else:
            yield shape


def analyze_slide(slide, slide_width: float = 13.333, slide_height: float = 7.5) -> SourceProfile:
    """Extract source composition facts without mutating the source slide."""

    profile = SourceProfile(slide_width, slide_height)
    for z_order, shape in enumerate(iter_source_shapes(slide.shapes)):
        bounds = _bounds(shape, slide_width, slide_height)
        shape_type = getattr(shape, "shape_type", None)
        if getattr(shape, "has_table", False):
            table = shape.table
            profile.elements.append(SourceElement(
                ElementKind.TABLE,
                bounds,
                row_count=len(table.rows),
                column_count=len(table.columns),
                z_order=z_order,
            ))
        elif shape_type == MSO_SHAPE_TYPE.PICTURE:
            profile.elements.append(SourceElement(
                ElementKind.IMAGE,
                bounds,
                image_blob=shape.image.blob,
                z_order=z_order,
            ))
        elif getattr(shape, "has_text_frame", False) and shape.text.strip():
            paragraphs = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            profile.elements.append(SourceElement(
                ElementKind.TEXT,
                bounds,
                text="\n".join(p.text.strip() for p in paragraphs),
                role=_text_role(shape, shape.text, slide_width, slide_height),
                paragraph_count=len(paragraphs),
                has_bullets=any(p.level > 0 for p in paragraphs),
                z_order=z_order,
            ))
        elif shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.AUTO_SHAPE):
            profile.elements.append(SourceElement(ElementKind.SHAPE, bounds, z_order=z_order))
    return profile


def choose_source_template(profile: SourceProfile, title: str, body: str) -> str:
    """Choose a composition family using source geometry before word heuristics."""

    upper = title.upper()
    if profile.table_elements:
        return "table_dashboard"
    if profile.image_elements and profile.text_density < 0.13:
        return "image_story" if len(profile.image_elements) == 1 else "image_collage"
    if profile.connector_count >= 2 and len(profile.text_elements) >= 3:
        return "flowchart"
    if profile.has_two_columns and profile.text_density > 0.22:
        return "comparison_grid"
    if "КЛЮЧЕВЫЕ СЛОВА" in upper:
        return "key_points"
    if upper.startswith(("ВОПРОСЫ", "ЗАДАНИЯ", "ЗАДАЧА")):
        return "exercise_card"
    if upper.startswith(("ПРИМЕР", "РЕШЕНИЕ")):
        return "solution_steps"
    if "ФОРМУЛ" in upper or "ЗАКОНОМЕРНОСТ" in upper:
        return "formula_focus"
    if len(body) > 600 or profile.text_density > 0.38:
        return "content_overview"
    return "auto"